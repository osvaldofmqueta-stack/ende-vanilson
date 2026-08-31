"""Generate the academic defense deck and its browser preview."""

from base64 import b64encode
from pathlib import Path
from textwrap import dedent

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "outputs"
STATIC_DIR = ROOT / "static"
OUT_DIR.mkdir(exist_ok=True)

PPTX_PATH = OUT_DIR / "apresentacao_defesa_ende_2026.pptx"
PREVIEW_PATH = STATIC_DIR / "presentation_preview.html"
LOGO_PATH = STATIC_DIR / "images" / "logo_iscat.png"

NAVY = "0B1F33"
NAVY_2 = "102D47"
BLUE = "2F76FF"
TEAL = "16B8A6"
MINT = "B8F1E7"
YELLOW = "F4C95D"
WHITE = "F7FAFC"
MUTED = "B8C7D5"
INK = "17324D"
PALE = "EAF2F6"
RED = "F47C7C"


def rgb(value):
    return RGBColor.from_string(value)


def set_background(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, lines, x, y, w, h, size=16, color=WHITE,
                  spacing=5, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(spacing)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        if bullet:
            p.text = "• " + line
    return box


def add_rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def add_circle(slide, x, y, d, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_title(slide, kicker, title, number):
    add_text(slide, kicker.upper(), 0.65, 0.36, 3.4, 0.28, 10, TEAL, True)
    add_text(slide, title, 0.65, 0.72, 11.4, 0.65, 26, WHITE, True)
    add_text(slide, f"{number:02d}", 12.2, 0.42, 0.5, 0.32, 11, MUTED, True, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0.65, 1.48, 12.0, 0.02, TEAL)


def add_footer(slide, label="Defesa de monografia • 2026"):
    add_text(slide, label, 0.65, 7.12, 4.8, 0.2, 8.5, MUTED)
    add_text(slide, "ENDE Central", 10.2, 7.12, 2.45, 0.2, 8.5, MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=TEAL, icon=None):
    add_rect(slide, x, y, w, h, NAVY_2, True, NAVY_2)
    add_rect(slide, x, y, 0.06, h, accent)
    if icon:
        add_circle(slide, x + 0.24, y + 0.25, 0.42, accent)
        add_text(slide, icon, x + 0.24, y + 0.27, 0.42, 0.34, 14, NAVY, True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        title_x = x + 0.82
    else:
        title_x = x + 0.26
    add_text(slide, title, title_x, y + 0.23, w - (title_x - x) - 0.2, 0.3, 14, WHITE, True)
    add_text(slide, body, x + 0.26, y + 0.72, w - 0.48, h - 0.9, 11.5, MUTED)


def add_pill(slide, text, x, y, w, color=TEAL, text_color=NAVY):
    add_rect(slide, x, y, w, 0.34, color, True, color)
    add_text(slide, text, x + 0.05, y + 0.04, w - 0.1, 0.23, 9.5, text_color, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_step(slide, x, number, title, body, accent=TEAL):
    add_circle(slide, x, 3.05, 0.56, accent)
    add_text(slide, str(number), x, 3.15, 0.56, 0.26, 14, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, title, x - 0.18, 3.78, 1.65, 0.3, 13, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, body, x - 0.36, 4.18, 2.0, 0.88, 10.5, MUTED, align=PP_ALIGN.CENTER)


def add_image_fit(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — cover
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_rect(slide, 8.82, 0, 4.52, 7.5, NAVY_2)
    add_rect(slide, 8.82, 0, 0.12, 7.5, TEAL)
    if LOGO_PATH.exists():
        add_image_fit(slide, LOGO_PATH, 9.52, 0.7, 2.2, 1.2)
    add_pill(slide, "DEFESA • 2026", 0.72, 0.7, 1.55, YELLOW)
    add_text(slide, "DESENVOLVIMENTO DE SISTEMA DE GESTÃO E CONTROLO DE PAGAMENTO DE ENERGIA PRÉ-PAGO E PÓS-PAGO", 0.72, 1.55, 7.5, 1.85, 28, WHITE, True)
    add_text(slide, "Caso de estudo realizado no Posto de Atendimento da ENDE Central", 0.76, 3.7, 6.9, 0.72, 17, MINT)
    add_rect(slide, 0.76, 5.15, 2.25, 0.03, TEAL)
    add_text(slide, "Adão Joveta da Costa", 0.76, 5.43, 3.8, 0.3, 14, WHITE, True)
    add_text(slide, "Orientador: Eng.º Osvaldo Fernando Muondo Queta", 0.76, 5.82, 5.7, 0.3, 11.5, MUTED)
    add_text(slide, "Licenciatura em Engenharia Informática\nInstituto Superior Politécnico Privado da Catepa", 9.52, 5.55, 2.95, 0.72, 11, WHITE, True)
    add_text(slide, "Gestão • Consumo • Pagamentos", 9.52, 6.78, 2.8, 0.25, 10, TEAL, True)

    # 2 — context
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "01 • Contexto", "O desafio que motivou o sistema", 2)
    add_text(slide, "No atendimento de energia, informação dispersa significa mais tempo para localizar dados, conferir consumos e acompanhar pagamentos.", 0.78, 1.84, 8.2, 0.58, 18, MINT)
    add_card(slide, 0.78, 3.0, 3.65, 2.55, "Processos dispersos", "Cadastro, contratos, contadores e pagamentos precisam de uma visão integrada para apoiar o atendimento.", BLUE, "1")
    add_card(slide, 4.84, 3.0, 3.65, 2.55, "Controlo manual", "A conferência de leituras, recargas e faturas pode gerar retrabalho e dificultar o acompanhamento.", TEAL, "2")
    add_card(slide, 8.9, 3.0, 3.65, 2.55, "Visibilidade limitada", "Sem relatórios organizados, torna-se mais difícil acompanhar consumo, receitas e dívidas.", YELLOW, "3")
    add_footer(slide)

    # 3 — objectives
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "02 • Direção", "Objetivo e pergunta de partida", 3)
    add_text(slide, "Como uma aplicação web pode apoiar o controlo do pagamento de energia pré-paga e pós-paga no Posto de Atendimento da ENDE Central?", 0.78, 1.82, 11.4, 0.72, 21, WHITE, True)
    add_rect(slide, 0.78, 3.0, 4.3, 2.55, TEAL, True, TEAL)
    add_text(slide, "OBJETIVO GERAL", 1.08, 3.32, 3.5, 0.25, 10, NAVY, True)
    add_text(slide, "Desenvolver uma aplicação web contextualizada para organizar dados e apoiar a gestão de energia pré-paga e pós-paga.", 1.08, 3.78, 3.55, 1.25, 17, NAVY, True)
    add_text(slide, "Objetivos específicos", 5.7, 3.0, 3.6, 0.3, 16, TEAL, True)
    add_rich_text(slide, [
        "Centralizar clientes, contratos, contadores e tarifas.",
        "Registar leituras, recargas, faturas e pagamentos.",
        "Disponibilizar consultas e relatórios para apoiar decisões."
    ], 5.68, 3.48, 6.2, 1.65, 15, WHITE, 12, True)
    add_footer(slide)

    # 4 — methodology
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "03 • Método", "Como o trabalho foi desenvolvido", 4)
    add_step(slide, 0.98, 1, "Explorar", "Pesquisa exploratória e descritiva sobre o problema.", BLUE)
    add_step(slide, 3.35, 2, "Compreender", "Revisão bibliográfica, análise documental e entrevistas.", TEAL)
    add_step(slide, 5.72, 3, "Modelar", "Requisitos, entidades, perfis e fluxos do MVP.", YELLOW)
    add_step(slide, 8.09, 4, "Construir", "Implementação modular com Django e base PostgreSQL.", BLUE)
    add_step(slide, 10.46, 5, "Verificar", "Testes funcionais e demonstração com dados simulados.", TEAL)
    add_rect(slide, 1.0, 5.75, 11.1, 0.68, NAVY_2, True, NAVY_2)
    add_text(slide, "Nota de rigor: os dados disponíveis das entrevistas são parciais; não se afirma saturação nem generalização estatística.", 1.28, 5.95, 10.55, 0.26, 12.5, MINT, True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 5 — roles
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "04 • Utilizadores", "Quatro perfis, responsabilidades diferentes", 5)
    roles = [
        ("Administrador", "Configura perfis e acompanha a gestão global.", BLUE, "A"),
        ("Operador", "Atende, cadastra clientes e regista operações.", TEAL, "O"),
        ("Financeiro", "Acompanha faturas, pagamentos, receitas e dívidas.", YELLOW, "F"),
        ("Cliente", "Consulta os seus dados, consumo, saldo e faturas.", MINT, "C"),
    ]
    for i, (title, body, color, icon) in enumerate(roles):
        x = 0.78 + (i % 2) * 6.1
        y = 2.12 + (i // 2) * 2.05
        add_rect(slide, x, y, 5.52, 1.5, NAVY_2, True, NAVY_2)
        add_circle(slide, x + 0.35, y + 0.38, 0.72, color)
        add_text(slide, icon, x + 0.35, y + 0.53, 0.72, 0.28, 17, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 1.35, y + 0.32, 3.7, 0.3, 16, WHITE, True)
        add_text(slide, body, x + 1.35, y + 0.76, 3.75, 0.4, 12, MUTED)
    add_footer(slide)

    # 6 — MVP
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "05 • MVP", "O que foi priorizado na aplicação", 6)
    modules = [
        ("Cadastro", "Clientes\nContratos\nTarifas", BLUE),
        ("Medição", "Contadores\nLeituras\nConsumo", TEAL),
        ("Financeiro", "Recargas\nFaturas\nPagamentos", YELLOW),
        ("Análise", "Relatórios\nPDF e Excel\nConsultas", MINT),
    ]
    for i, (title, body, color) in enumerate(modules):
        x = 0.78 + i * 3.05
        add_rect(slide, x, 2.15, 2.55, 2.45, NAVY_2, True, NAVY_2)
        add_rect(slide, x, 2.15, 2.55, 0.1, color)
        add_text(slide, title, x + 0.25, 2.58, 2.05, 0.3, 16, WHITE, True)
        add_text(slide, body, x + 0.25, 3.18, 2.05, 1.0, 14, color, True)
        if i < 3:
            add_text(slide, "→", x + 2.63, 3.04, 0.4, 0.4, 22, MUTED, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.5, 5.2, 10.3, 0.78, NAVY_2, True, NAVY_2)
    add_text(slide, "Recibos e notificações foram modelados, mas os fluxos completos de emissão/envio permanecem fora do MVP.", 1.8, 5.45, 9.7, 0.26, 12.5, YELLOW, True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 7 — architecture
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "06 • Solução", "Arquitetura tecnológica", 7)
    add_rect(slide, 0.82, 2.0, 3.35, 3.35, NAVY_2, True, NAVY_2)
    add_text(slide, "INTERFACE", 1.15, 2.35, 2.7, 0.25, 10, TEAL, True)
    add_text(slide, "Bootstrap 5\nHTML responsivo\nMensagens em português", 1.15, 2.78, 2.7, 1.05, 17, WHITE, True)
    add_rect(slide, 4.98, 2.0, 3.35, 3.35, NAVY_2, True, NAVY_2)
    add_text(slide, "APLICAÇÃO", 5.31, 2.35, 2.7, 0.25, 10, BLUE, True)
    add_text(slide, "Python 3.11\nDjango 5.2\nDRF e django-filter", 5.31, 2.78, 2.7, 1.05, 17, WHITE, True)
    add_rect(slide, 9.14, 2.0, 3.35, 3.35, NAVY_2, True, NAVY_2)
    add_text(slide, "DADOS E SAÍDA", 9.47, 2.35, 2.7, 0.25, 10, YELLOW, True)
    add_text(slide, "PostgreSQL\nReportLab para PDF\nOpenPyXL para Excel", 9.47, 2.78, 2.7, 1.05, 17, WHITE, True)
    for x in (4.33, 8.48):
        add_text(slide, "→", x, 3.25, 0.48, 0.45, 24, TEAL, True, align=PP_ALIGN.CENTER)
    add_text(slide, "A configuração oficial do ambiente utiliza PostgreSQL; instalações locais sem servidor configurado usam SQLite apenas como fallback.", 1.0, 5.85, 11.2, 0.42, 11.5, MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 8 — data model
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "07 • Modelo", "Do cliente ao pagamento: informação ligada", 8)
    entities = [
        ("Cliente", "perfil • contrato", BLUE),
        ("Contador", "leituras • consumo", TEAL),
        ("Recarga", "saldo pré-pago", YELLOW),
        ("Fatura", "consumo • vencimento", MINT),
        ("Pagamento", "liquidação • estado", BLUE),
    ]
    for i, (title, body, color) in enumerate(entities):
        x = 0.72 + i * 2.48
        add_rect(slide, x, 2.45, 2.05, 1.55, NAVY_2, True, NAVY_2)
        add_rect(slide, x, 2.45, 2.05, 0.08, color)
        add_text(slide, title, x + 0.18, 2.83, 1.68, 0.28, 14, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.18, 3.28, 1.68, 0.34, 10.5, MUTED, align=PP_ALIGN.CENTER)
        if i < len(entities) - 1:
            add_text(slide, "→", x + 2.08, 3.0, 0.37, 0.35, 19, TEAL, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.35, 4.75, 10.6, 0.92, TEAL, True, TEAL)
    add_text(slide, "A modelação preserva relações entre cadastro, medição e financeiro para apoiar consultas e relatórios.", 1.72, 5.08, 9.9, 0.3, 14, NAVY, True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 9 — demo
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "08 • Demonstração", "Um atendimento acompanhado de ponta a ponta", 9)
    steps = [
        ("1", "Cadastrar", "Cliente, contrato e tarifa"),
        ("2", "Medir", "Contador, leitura e consumo"),
        ("3", "Registar", "Recarga ou fatura"),
        ("4", "Acompanhar", "Pagamento, dívida e relatório"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.9 + i * 3.05
        add_circle(slide, x + 0.84, 2.16, 0.68, [BLUE, TEAL, YELLOW, MINT][i])
        add_text(slide, num, x + 0.84, 2.32, 0.68, 0.26, 16, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x, 3.2, 2.36, 0.3, 16, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.12, 3.72, 2.12, 0.58, 12, MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, "→", x + 2.48, 2.36, 0.4, 0.35, 21, TEAL, True, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.02, 5.2, 9.3, 0.75, NAVY_2, True, NAVY_2)
    add_text(slide, "A demonstração utiliza dados simulados para validar o percurso funcional.", 2.28, 5.45, 8.8, 0.26, 13, YELLOW, True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 10 — results
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "09 • Resultados", "O que a avaliação funcional mostrou", 10)
    add_pill(slide, "DADOS SIMULADOS", 10.35, 1.86, 2.1, YELLOW)
    result_cards = [
        ("Mais centralização", "As informações principais passam a estar reunidas num único sistema.", BLUE),
        ("Cálculos apoiados", "Consumo, saldo, valores de fatura e pagamentos têm regras no sistema.", TEAL),
        ("Relatórios úteis", "A consulta e a exportação apoiam o acompanhamento financeiro e operacional.", YELLOW),
    ]
    for i, (title, body, color) in enumerate(result_cards):
        add_card(slide, 0.82 + i * 4.08, 2.55, 3.7, 2.35, title, body, color, str(i + 1))
    add_text(slide, "Estes resultados indicam potencial de melhoria do atendimento, mas não constituem medição de desempenho real da ENDE.", 1.0, 5.65, 11.15, 0.35, 12, MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide)

    # 11 — limitations
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "10 • Honestidade", "Limitações e próximos passos", 11)
    add_text(slide, "O MVP demonstra a base do sistema. A defesa deve apresentar claramente o que ainda requer evolução.", 0.8, 1.85, 10.5, 0.45, 17, MINT)
    add_rich_text(slide, [
        "Integrações reais com Multicaixa, ATM, USSD ou outras plataformas.",
        "Fluxo completo de emissão de recibos e envio de notificações.",
        "Refinamento das permissões e cobertura de testes por perfil.",
        "Automatização da suspensão por dívida e suporte mais amplo ao pós-pago."
    ], 0.9, 2.75, 6.4, 2.45, 15, WHITE, 15, True)
    add_rect(slide, 8.15, 2.55, 3.85, 2.75, NAVY_2, True, NAVY_2)
    add_text(slide, "EVOLUÇÃO PROPOSTA", 8.52, 2.92, 3.1, 0.25, 10, TEAL, True)
    add_text(slide, "Consolidar entrevistas, ampliar testes e validar o sistema com dados operacionais autorizados.", 8.52, 3.45, 2.95, 1.05, 17, WHITE, True)
    add_footer(slide)

    # 12 — conclusion
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_rect(slide, 0, 0, 0.16, 7.5, TEAL)
    add_text(slide, "11 • CONCLUSÃO", 0.85, 0.75, 3.5, 0.3, 11, TEAL, True)
    add_text(slide, "Uma base integrada para\num atendimento mais organizado.", 0.85, 1.5, 7.8, 1.35, 31, WHITE, True)
    add_text(slide, "O sistema reúne o cadastro, a medição, o pré-pago, o pós-pago, os pagamentos e os relatórios num MVP funcional, com PostgreSQL no ambiente oficial e dados simulados na demonstração.", 0.9, 3.35, 7.0, 1.05, 17, MINT)
    add_rect(slide, 9.0, 1.55, 2.65, 2.65, TEAL, True, TEAL)
    add_text(slide, "MVP\n+\nVALIDAÇÃO\nFUNCIONAL", 9.22, 2.05, 2.2, 1.65, 20, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Obrigado.\nPerguntas?", 0.9, 5.65, 4.5, 0.72, 24, YELLOW, True)
    add_text(slide, "Adão Joveta da Costa • 2026", 9.0, 6.7, 3.3, 0.25, 10, MUTED, align=PP_ALIGN.RIGHT)

    prs.save(PPTX_PATH)


SLIDES = [
    ("Capa", "Sistema de Gestão de Energia Pré-paga e Pós-paga", "ENDE Central • Defesa 2026"),
    ("O desafio", "Processos dispersos, controlo manual e visibilidade limitada.", "Contexto do atendimento"),
    ("Objetivo", "Organizar dados e apoiar a gestão de energia pré-paga e pós-paga.", "Pergunta de partida"),
    ("Metodologia", "Explorar → compreender → modelar → construir → verificar.", "Entrevistas com dados parciais"),
    ("Utilizadores", "Administrador • Operador • Financeiro • Cliente", "Perfis oficiais"),
    ("MVP", "Cadastro • medição • financeiro • análise", "Recibos e notificações modelados, ainda parciais"),
    ("Solução", "Django + PostgreSQL + Bootstrap + relatórios", "Arquitetura tecnológica"),
    ("Modelo", "Cliente → contador → recarga/fatura → pagamento", "Informação ligada"),
    ("Demonstração", "Cadastrar → medir → registar → acompanhar", "Dados simulados"),
    ("Resultados", "Centralização, cálculos apoiados e relatórios úteis.", "Avaliação funcional"),
    ("Limitações", "Integrações, notificações, permissões, testes e automação.", "Próximos passos"),
    ("Conclusão", "Uma base integrada para um atendimento mais organizado.", "Obrigado • Perguntas"),
]


def image_data(path):
    if not path.exists():
        return ""
    mime = "image/png" if path.suffix.lower() == ".png" else "image/jpeg"
    return f"data:{mime};base64,{b64encode(path.read_bytes()).decode('ascii')}"


def build_preview():
    logo = image_data(LOGO_PATH)
    cards = []
    for index, (kicker, title, subtitle) in enumerate(SLIDES):
        cards.append(
            f"""
            <section class="slide {'active' if index == 0 else ''}" data-index="{index}">
              <div class="slide-number">{index + 1:02d}</div>
              <div class="accent"></div>
              <div class="kicker">{kicker.upper()}</div>
              <h1>{title}</h1>
              <p class="subtitle">{subtitle}</p>
              <div class="visual visual-{index + 1}">
                <span></span><span></span><span></span>
              </div>
              <div class="meta">Defesa de monografia <b>•</b> 2026 <em>ENDE Central</em></div>
            </section>
            """
        )
    html = f"""<!doctype html>
<html lang="pt">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Apresentação da defesa — ENDE</title>
<style>
:root {{ --navy:#0B1F33; --navy2:#102D47; --teal:#16B8A6; --mint:#B8F1E7; --yellow:#F4C95D; --muted:#B8C7D5; --blue:#2F76FF; }}
* {{ box-sizing:border-box; }}
body {{ margin:0; min-height:100vh; background:#071522; color:#F7FAFC; font-family:Inter,Segoe UI,Arial,sans-serif; display:grid; place-items:center; overflow:hidden; }}
.deck {{ width:min(94vw,1200px); aspect-ratio:16/9; position:relative; border-radius:18px; overflow:hidden; box-shadow:0 24px 80px #0008; }}
.slide {{ display:none; position:absolute; inset:0; padding:6.5% 7%; background:var(--navy); overflow:hidden; }}
.slide.active {{ display:block; }}
.slide:after {{ content:""; position:absolute; right:-12%; top:-25%; width:48%; height:150%; background:var(--navy2); transform:rotate(8deg); opacity:.9; }}
.slide-number {{ position:absolute; right:6%; top:7%; color:var(--muted); font-weight:700; font-size:clamp(10px,1.3vw,16px); z-index:2; }}
.accent {{ height:3px; width:92%; background:var(--teal); position:absolute; top:21%; left:7%; z-index:2; }}
.kicker {{ color:var(--teal); font-weight:800; letter-spacing:.14em; font-size:clamp(9px,1.25vw,15px); position:relative; z-index:2; }}
h1 {{ font-size:clamp(24px,4vw,56px); line-height:1.05; max-width:74%; margin:4% 0 2%; position:relative; z-index:2; }}
.subtitle {{ color:var(--mint); font-size:clamp(13px,1.7vw,23px); max-width:64%; position:relative; z-index:2; }}
.visual {{ position:absolute; right:9%; top:34%; width:27%; height:32%; display:flex; flex-direction:column; gap:10%; z-index:2; }}
.visual span {{ display:block; height:22%; border-radius:10px; background:var(--teal); opacity:.9; }}
.visual span:nth-child(2) {{ width:78%; background:var(--yellow); }}
.visual span:nth-child(3) {{ width:56%; background:var(--blue); }}
.visual-4 span:nth-child(1),.visual-8 span:nth-child(1) {{ background:var(--blue); }}
.visual-5 span:nth-child(2),.visual-9 span:nth-child(2) {{ background:var(--mint); }}
.visual-10 span:nth-child(3),.visual-11 span:nth-child(3) {{ background:var(--yellow); }}
.meta {{ position:absolute; bottom:5%; left:7%; color:var(--muted); font-size:clamp(9px,1.1vw,14px); z-index:2; }}
.meta em {{ position:absolute; left:clamp(190px,55vw,700px); width:180px; font-style:normal; text-align:right; }}
.controls {{ position:absolute; bottom:2.5%; left:50%; transform:translateX(-50%); display:flex; gap:8px; z-index:5; }}
button {{ border:1px solid #ffffff33; background:#ffffff18; color:white; border-radius:999px; padding:7px 14px; cursor:pointer; }}
button:hover {{ background:var(--teal); color:var(--navy); }}
.counter {{ position:absolute; right:3%; bottom:3%; color:#ffffff99; font-size:12px; z-index:5; }}
@media(max-width:700px) {{ .slide {{ padding:8% 8%; }} .accent {{ top:24%; }} .visual {{ opacity:.55; right:4%; width:34%; }} h1 {{ max-width:90%; }} .subtitle {{ max-width:84%; }} .meta em {{ display:none; }} }}
</style>
</head>
<body>
<div class="deck" aria-label="Prévia da apresentação">
  {''.join(cards)}
  <div class="controls"><button id="prev">Anterior</button><button id="next">Seguinte</button></div>
  <div class="counter" id="counter">01 / {len(SLIDES):02d}</div>
</div>
<script>
const slides=[...document.querySelectorAll('.slide')]; let current=0;
function show(n) {{ current=(n+slides.length)%slides.length; slides.forEach((s,i)=>s.classList.toggle('active',i===current)); document.querySelector('#counter').textContent=String(current+1).padStart(2,'0')+' / '+String(slides.length).padStart(2,'0'); }}
document.querySelector('#prev').onclick=()=>show(current-1); document.querySelector('#next').onclick=()=>show(current+1);
document.addEventListener('keydown',e=>{{if(e.key==='ArrowRight'||e.key===' ')show(current+1);if(e.key==='ArrowLeft')show(current-1);}});
</script>
</body></html>"""
    PREVIEW_PATH.write_text(html, encoding="utf-8")


if __name__ == "__main__":
    build_pptx()
    build_preview()
    print(PPTX_PATH)
    print(PREVIEW_PATH)